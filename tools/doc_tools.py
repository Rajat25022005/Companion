import json
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

WORKSPACE_DIR = Path(__file__).parent.parent / 'workspace'
WORKSPACE_DIR.mkdir(exist_ok=True)


def _resolve_to_workspace(filename: str) -> Path:
    """Ensure output files always land in the workspace directory."""
    p = Path(filename)
    if not p.is_absolute():
        p = WORKSPACE_DIR / p.name
    return p


def _sanitize_latin1(text: str) -> str:
    """Replace Unicode characters that Helvetica can't render with ASCII equivalents."""
    replacements = {
        '\u2192': '->', '\u2190': '<-', '\u2194': '<->', '\u21d2': '=>',
        '\u2022': '*', '\u2013': '-', '\u2014': '--', '\u2018': "'", '\u2019': "'",
        '\u201c': '"', '\u201d': '"', '\u2026': '...', '\u00b2': '2', '\u2082': '2',
        '\u2083': '3', '\u2084': '4', '\u2085': '5', '\u2086': '6',
        '\u00b0': 'deg', '\u00d7': 'x', '\u2248': '~=', '\u2260': '!=',
        '\u2264': '<=', '\u2265': '>=', '\u00b1': '+/-',
    }
    for k, v in replacements.items():
        text = text.replace(k, v)
    try:
        text.encode('latin-1')
    except UnicodeEncodeError:
        text = text.encode('latin-1', errors='replace').decode('latin-1')
    return text


def create_pdf(content: str, filename: str, title: str = '', author: str = '') -> dict:
    """Generate a PDF document from text content and save it to workspace."""
    try:
        from fpdf import FPDF

        content = _sanitize_latin1(content)
        title = _sanitize_latin1(title) if title else ''
        author = _sanitize_latin1(author) if author else ''

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        if title:
            pdf.set_font('Helvetica', 'B', 20)
            pdf.cell(0, 15, title, ln=True, align='C')
            if author:
                pdf.set_font('Helvetica', 'I', 12)
                pdf.cell(0, 8, author, ln=True, align='C')
            pdf.ln(10)

        pdf.set_font('Helvetica', '', 11)
        for line in content.split('\n'):
            stripped = line.strip()
            if stripped.startswith('# '):
                pdf.ln(6)
                pdf.set_font('Helvetica', 'B', 16)
                pdf.cell(0, 10, stripped[2:], ln=True)
                pdf.set_font('Helvetica', '', 11)
            elif stripped.startswith('## '):
                pdf.ln(4)
                pdf.set_font('Helvetica', 'B', 13)
                pdf.cell(0, 8, stripped[3:], ln=True)
                pdf.set_font('Helvetica', '', 11)
            elif stripped.startswith('- '):
                pdf.cell(10)
                pdf.multi_cell(0, 6, f'\u2022 {stripped[2:]}')
            elif stripped:
                pdf.multi_cell(0, 6, stripped)
            else:
                pdf.ln(3)

        out = _resolve_to_workspace(filename)
        pdf.output(str(out))
        return {
            'status': 'ok',
            'path': str(out),
            'filename': out.name,
            'download_url': f'/files/{out.name}',
            'pages': pdf.page,
        }
    except ImportError:
        return {'error': 'fpdf2 not installed. Run: pip install fpdf2'}
    except Exception as e:
        return {'error': str(e)}

create_pdf._tool_schema = {
    'parameters': {
        'type': 'object',
        'properties': {
            'content': {'type': 'string', 'description': 'Text content for the PDF'},
            'filename': {'type': 'string', 'description': 'Output filename (e.g. report.pdf)'},
            'title': {'type': 'string', 'description': 'Document title'},
            'author': {'type': 'string', 'description': 'Author name'},
        },
        'required': ['content', 'filename'],
    }
}


def create_docx(content: str, filename: str, title: str = '') -> dict:
    """Generate a Word document from text content and save it to workspace."""
    try:
        from docx import Document
        from docx.shared import Pt

        doc = Document()
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)

        if title:
            doc.add_heading(title, level=0)

        for line in content.split('\n'):
            stripped = line.strip()
            if stripped.startswith('# '):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith('## '):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith('### '):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith('- '):
                doc.add_paragraph(stripped[2:], style='List Bullet')
            elif stripped:
                doc.add_paragraph(stripped)

        out = _resolve_to_workspace(filename)
        doc.save(str(out))
        return {
            'status': 'ok',
            'path': str(out),
            'filename': out.name,
            'download_url': f'/files/{out.name}',
            'paragraphs': len(doc.paragraphs),
        }
    except ImportError:
        return {'error': 'python-docx not installed. Run: pip install python-docx'}
    except Exception as e:
        return {'error': str(e)}

create_docx._tool_schema = {
    'parameters': {
        'type': 'object',
        'properties': {
            'content': {'type': 'string', 'description': 'Text content'},
            'filename': {'type': 'string', 'description': 'Output filename (e.g. report.docx)'},
            'title': {'type': 'string', 'description': 'Document title'},
        },
        'required': ['content', 'filename'],
    }
}


def create_pptx(slides: str, filename: str, title: str = '') -> dict:
    """Generate a PowerPoint presentation and save it to workspace. Slides separated by '---'."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        if title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title

        for slide_text in [s.strip() for s in slides.split('---') if s.strip()]:
            lines = slide_text.split('\n')
            slide_title = lines[0].lstrip('# ').strip() if lines else 'Slide'
            bullets = [l.lstrip('- ').strip() for l in lines[1:] if l.strip()]

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title

            if bullets and slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet

        out = _resolve_to_workspace(filename)
        prs.save(str(out))
        return {
            'status': 'ok',
            'path': str(out),
            'filename': out.name,
            'download_url': f'/files/{out.name}',
            'slides': len(prs.slides),
        }
    except ImportError:
        return {'error': 'python-pptx not installed. Run: pip install python-pptx'}
    except Exception as e:
        return {'error': str(e)}

create_pptx._tool_schema = {
    'parameters': {
        'type': 'object',
        'properties': {
            'slides': {'type': 'string', 'description': 'Slide content separated by ---'},
            'filename': {'type': 'string', 'description': 'Output filename (e.g. deck.pptx)'},
            'title': {'type': 'string', 'description': 'Title slide text'},
        },
        'required': ['slides', 'filename'],
    }
}

DOC_TOOLS = {
    'create_pdf': create_pdf,
    'create_docx': create_docx,
    'create_pptx': create_pptx,
}
